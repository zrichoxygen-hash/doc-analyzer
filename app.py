import streamlit as st
import os
from pathlib import Path
from openai import OpenAI
from dotenv import load_dotenv
from tqdm import tqdm
import pandas as pd
import re
import time
import json
import io
import base64
from PyPDF2 import PdfReader
from pptx import Presentation
import pypdfium2 as pdfium

# Configuration de la page
st.set_page_config(page_title="Document Analyzer", layout="wide")

# Charger les variables d'environnement
load_dotenv("API.env")

# Initialiser le client OpenAI
# En priorité: variables d'environnement (Render), sinon Streamlit Secrets (local dev)
api_key = os.getenv("OPENAI_API_KEY")

if not api_key:
    try:
        # Fallback sur Streamlit Secrets pour dev local avec secrets.toml
        api_key = st.secrets.get("OPENAI_API_KEY")
    except:
        pass

if not api_key:
    st.error("❌ Clé API OpenAI non trouvée! Configurez OPENAI_API_KEY en variable d'environnement (Render) ou dans Secrets Streamlit/API.env (local)")
    st.stop()

client = OpenAI(api_key=api_key)

# ===== FONCTION DE PERSISTANCE DES CRITÈRES =====
CRITERIA_FILE = "saved_criteria.json"

def load_criteria():
    """Charge les critères sauvegardés ou retourne les critères par défaut."""
    if Path(CRITERIA_FILE).exists():
        try:
            with open(CRITERIA_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                return data.get("criteria", get_default_criteria())
        except:
            return get_default_criteria()
    return get_default_criteria()

def get_default_criteria():
    """Retourne les critères par défaut."""
    return """Veuillez évaluer ce document académique sur la base des 5 critères suivants et attribuer une note totale sur 20 points (4 points par critère).

Les critères sont:
1. Plan adopté (4 pts) - Cohérence et organisation logique du contenu
2. Design des slides (4 pts) - Esthétique, lisibilité et qualité de la présentation visuelle
3. Référencement (4 pts) - Utilisation correcte des normes (APA ou autre) pour les sources
4. Existence de chiffres, tableaux ou graphiques (4 pts) - Présence et pertinence des éléments visuels de données
5. Richesse des informations (4 pts) - Profondeur du contenu, vocabulaire technique approprié et attesté par les références

Veuillez fournir:
- Une analyse succincte pour chaque critère (1-2 lignes max)
- La note obtenue pour chaque critère (0-4)
- La note totale sur 20"""

def save_criteria(criteria_text):
    """Sauvegarde les critères modifiés dans un fichier JSON."""
    with open(CRITERIA_FILE, "w", encoding="utf-8") as f:
        json.dump({"criteria": criteria_text}, f, ensure_ascii=False, indent=2)

# ===== DÉFINITION DES FONCTIONS =====

def parse_filename(filename):
    """Extrait le prénom et nom du format: 'DESCRIPTION - FirstName LASTNAME.extension'"""
    try:
        name_without_ext = Path(filename).stem
        
        if " - " in name_without_ext:
            parts = name_without_ext.split(" - ")
            if len(parts) >= 2:
                names = " - ".join(parts[1:]).strip()
                name_parts = names.split()
                
                if len(name_parts) >= 2:
                    first_name = name_parts[0]
                    last_name = " ".join(name_parts[1:])
                    return first_name, last_name
                elif len(name_parts) == 1:
                    return name_parts[0], ""
    except Exception as e:
        pass
    
    return None, None


def extract_note(response_text):
    """Extrait la note totale du texte de réponse avec plusieurs patterns possibles."""
    try:
        # Pattern 1: "NOTE TOTALE: X/20"
        match = re.search(r'NOTE\s+TOTALE\s*:\s*(\d+(?:[.,]\d+)?)\s*/\s*20', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 2: "NOTE TOTALE: X" (sans /20)
        match = re.search(r'NOTE\s+TOTALE\s*:\s*(\d+(?:[.,]\d+)?)', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 3: Chercher juste "/20" suivi d'un nombre
        match = re.search(r'(\d+(?:[.,]\d+)?)\s*/\s*20', response_text)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 4: Chercher "note:" ou "note =" suivi d'un nombre
        match = re.search(r'note\s*[:=]\s*(\d+(?:[.,]\d+)?)', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
    except Exception as e:
        pass
    
    return None


def read_upload_results(filename="upload_results.txt"):
    """Lit le fichier de résultats d'upload pour extraire les IDs de fichiers."""
    results = []
    
    try:
        with open(filename, "r", encoding="utf-8") as f:
            lines = f.readlines()
        
        current_file = None
        for line in lines:
            if line.startswith("Fichier:"):
                current_file = line.replace("Fichier:", "").strip()
            elif line.startswith("ID:") and current_file:
                file_id = line.replace("ID:", "").strip()
                results.append((file_id, current_file))
                current_file = None
        
        return results
    
    except FileNotFoundError:
        return []
    except Exception as e:
        return []


def ocr_pdf_to_text(file_path):
    """Effectue un OCR page par page d'un PDF en utilisant GPT-4o-mini Vision."""
    ocr_pages = []
    pdf_document = pdfium.PdfDocument(file_path)

    for page_index in range(len(pdf_document)):
        page = pdf_document[page_index]

        # Rendu en image pour OCR (qualite suffisante pour texte de documents académiques)
        bitmap = page.render(scale=2.0)
        pil_image = bitmap.to_pil()

        image_buffer = io.BytesIO()
        pil_image.save(image_buffer, format="JPEG", quality=90)
        image_b64 = base64.b64encode(image_buffer.getvalue()).decode("ascii")

        ocr_response = client.responses.create(
            model="gpt-4o-mini",
            input=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "input_text",
                            "text": "Fais un OCR fidele de cette page. Retourne uniquement le texte brut, sans commentaire."
                        },
                        {
                            "type": "input_image",
                            "image_url": f"data:image/jpeg;base64,{image_b64}"
                        }
                    ]
                }
            ],
            max_output_tokens=4000
        )

        page_text = (ocr_response.output_text or "").strip()
        ocr_pages.append(f"--- Page {page_index + 1} ---\n{page_text}")

    return "\n\n".join(ocr_pages)


def evaluate_document(file_id, file_name, evaluation_prompt, documents_folder=None):
    """Évalue un document en utilisant l'API OpenAI avec GPT-4 Vision."""
    first_name, last_name = parse_filename(file_name)
    
    # Dossier par défaut
    if documents_folder is None:
        documents_folder = r"C:\Users\hp\Downloads\Documents a analyser"
    
    try:
        file_path = os.path.join(documents_folder, file_name)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Fichier non trouvé: {file_path}")
        
        file_ext = Path(file_name).suffix.lower()
        extracted_text = ""
        
        if file_ext == ".pdf":
            try:
                # OCR complet de toutes les pages du PDF
                extracted_text = ocr_pdf_to_text(file_path)

                # Fallback texte natif si l'OCR ne retourne rien
                if not extracted_text.strip():
                    pdf_reader = PdfReader(file_path)
                    fallback_text_parts = []
                    for page_num, page in enumerate(pdf_reader.pages, start=1):
                        fallback_text_parts.append(f"--- Page {page_num} ---\n{page.extract_text() or ''}")
                    extracted_text = "\n\n".join(fallback_text_parts)
            except Exception as pdf_error:
                extracted_text = f"OCR indisponible pour {file_name}. Erreur: {str(pdf_error)}"
        
        elif file_ext == ".pptx":
            try:
                prs = Presentation(file_path)
                text_content = "Contenu de la présentation:\n"
                slide_count = len(prs.slides)
                
                for slide_idx in range(slide_count):
                    try:
                        slide = prs.slides[slide_idx]
                        text_content += f"\n--- Slide {slide_idx + 1} ---\n"
                        
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, "text") and shape.text:
                                    text_content += shape.text + "\n"
                            except:
                                pass
                    except:
                        continue

                extracted_text = text_content
            except Exception as pptx_error:
                extracted_text = f"Document PPTX: {file_name}. Erreur d'extraction: {str(pptx_error)}"

        else:
            extracted_text = f"Format non supporte pour extraction: {file_name}"

        user_prompt = (
            f"{evaluation_prompt}\n\n"
            f"Utilise le texte integral ci-dessous pour evaluer le document. "
            f"Retourne obligatoirement une NOTE TOTALE sur 20.\n\n"
            f"TEXTE INTEGRAL DU DOCUMENT:\n{extracted_text}"
        )
        
        response = client.chat.completions.create(
            model="gpt-4o",
            max_tokens=1024,
            messages=[
                {
                    "role": "system",
                    "content": "Tu es un evaluateur academique. Reponds avec une evaluation structuree et une ligne finale 'NOTE TOTALE: X/20'."
                },
                {
                    "role": "user",
                    "content": user_prompt
                }
            ]
        )
        
        evaluation = response.choices[0].message.content
        note = extract_note(evaluation)
        
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": note,
            "evaluation": evaluation,
            "status": "✅ Succès" if note is not None else "⚠️ Note non trouvée"
        }
    
    except FileNotFoundError as e:
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": None,
            "evaluation": str(e),
            "status": f"❌ Fichier non trouvé"
        }
    
    except Exception as e:
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": None,
            "evaluation": str(e),
            "status": f"❌ Erreur: {str(e)[:50]}"
        }


def upload_documents(folder_path):
    """Upload tous les documents PDF et PPTX d'un dossier vers l'API OpenAI.
    Le chemin peut être celui d'un dossier local (exécution locale) ou d'un
    répertoire temporaire créé par `st.file_uploader`.
    """
    ALLOWED_EXTENSIONS = {".pdf", ".pptx"}

    # Sauvegarder le dossier dans la session
    st.session_state.upload_folder = folder_path

    # Vérifier que le dossier existe
    if not os.path.exists(folder_path):
        st.error(f"❌ Le dossier n'existe pas: {folder_path}")
        return

    # Récupérer les fichiers
    documents = []
    try:
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)

            if os.path.isfile(file_path):
                file_ext = Path(file).suffix.lower()

                if file_ext in ALLOWED_EXTENSIONS:
                    documents.append(file_path)

        documents = sorted(documents)

    except Exception as e:
        st.error(f"❌ Erreur lors de la lecture du dossier: {e}")
        return

    if not documents:
        st.warning(f"⚠️ Aucun fichier PDF ou PPTX trouvé dans: {folder_path}")
        return

    st.info(f"📁 {len(documents)} fichier(s) à uploader")

    # Uploader les documents
    progress_bar = st.progress(0)
    status_text = st.empty()

    results = []
    for idx, file_path in enumerate(documents):
        # Mettre à jour la progress bar
        progress = (idx + 1) / len(documents)
        progress_bar.progress(progress)

        file_name = os.path.basename(file_path)
        status_text.text(f"Upload: {file_name}...")

        try:
            with open(file_path, "rb") as f:
                response = client.files.create(
                    file=(file_name, f),
                    purpose="assistants"
                )

            results.append({
                "file_name": file_name,
                "file_id": response.id,
                "status": "✅ Succès"
            })

            st.success(f"✅ {file_name}")

        except Exception as e:
            results.append({
                "file_name": file_name,
                "file_id": None,
                "status": f"❌ Erreur"
            })
            st.error(f"❌ {file_name}: {str(e)[:50]}")

    progress_bar.empty()
    status_text.empty()

    # Sauvegarder les résultats
    try:
        with open("upload_results.txt", "w", encoding="utf-8") as f:
            f.write("RÉSULTATS DE L'UPLOAD EN BATCH\n")
            f.write("=" * 60 + "\n\n")

            for result in results:
                f.write(f"Fichier: {result['file_name']}\n")
                f.write(f"Statut: {result['status']}\n")
                if result['file_id']:
                    f.write(f"ID: {result['file_id']}\n")
                f.write("\n")

        st.success(f"💾 Résultats sauvegardés dans upload_results.txt")

    except Exception as e:
        st.error(f"⚠️ Erreur lors de la sauvegarde: {e}")


def evaluate_documents(evaluation_prompt):
    """Évalue tous les documents en batch."""
    
    # Lire les fichiers uploadés
    files_to_evaluate = read_upload_results()
    
    if not files_to_evaluate:
        st.error("Aucun fichier à évaluer. Vérifiez upload_results.txt")
        return
    
    st.info(f"Évaluation de {len(files_to_evaluate)} document(s)...")
    
    # Progress bar
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    results = []
    for idx, (file_id, file_name) in enumerate(files_to_evaluate):
        # Mettre à jour la progress bar
        progress = (idx + 1) / len(files_to_evaluate)
        progress_bar.progress(progress)
        status_text.text(f"Évaluation: {file_name}...")
        
        # Évaluer le document avec le dossier stocké
        result = evaluate_document(file_id, file_name, evaluation_prompt, st.session_state.upload_folder)
        results.append(result)
    
    # Résultats
    st.session_state.results = results
    progress_bar.empty()
    status_text.empty()
    
    # Afficher les résultats
    successful = sum(1 for r in results if r['note'] is not None)
    st.success(f"✅ Évaluation terminée! {successful}/{len(results)} réussis")


# ===== INTERFACE MAIN =====

# Initialiser l'état de session
if "criteria" not in st.session_state:
    st.session_state.criteria = load_criteria()

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "results" not in st.session_state:
    st.session_state.results = []

if "upload_folder" not in st.session_state:
    st.session_state.upload_folder = r"C:\Users\hp\Downloads\Documents a analyser"

# En-tête
st.title("📊 Document Analyzer - Evaluation System")
st.markdown("---")

# ===== SECTION UPLOAD =====
st.subheader("📤 Upload des Documents")

st.markdown("*Vous pouvez soit sélectionner plusieurs fichiers PDF/PPTX, soit téléverser un **dossier compressé (.zip)** contenant tous les documents.*")

# uploader multiple files
uploaded_files = st.file_uploader(
    "Choisissez des fichiers PDF/PPTX ou un zip de dossier :",
    type=["pdf", "pptx", "zip"],
    accept_multiple_files=True
)

if uploaded_files:
    if st.button("📤 Uploader et traiter"):
        with st.spinner("Préparation des fichiers..."):
            temp_dir = "./_tmp_uploads"
            os.makedirs(temp_dir, exist_ok=True)

            # vider le répertoire temporaire pour éviter mélanges
            for existing in os.listdir(temp_dir):
                path = os.path.join(temp_dir, existing)
                if os.path.isfile(path):
                    os.remove(path)
                else:
                    import shutil
                    shutil.rmtree(path)

            # sauvegarder les uploads et décompresser si nécessaire
            import zipfile
            for up in uploaded_files:
                dest_path = os.path.join(temp_dir, up.name)
                with open(dest_path, "wb") as f:
                    f.write(up.getbuffer())

                if up.name.lower().endswith(".zip"):
                    try:
                        with zipfile.ZipFile(dest_path, "r") as z:
                            z.extractall(temp_dir)
                    except Exception as e:
                        st.error(f"Erreur lors de l'extraction du zip: {e}")
            # lancer l'upload sur le contenu du dossier temporaire
            upload_documents(temp_dir)

st.markdown("---")

# ===== SECTION ÉVALUATION =====

# Layout en deux colonnes
col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("🗨️ Modification des Critères (Optionnel)")
    
    # Chat interface
    chat_container = st.container()
    
    # Afficher l'historique du chat
    with chat_container:
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
    
    # Entrée utilisateur
    user_input = st.chat_input("Modifiez les critères d'évaluation...")
    
    if user_input:
        # Ajouter le message utilisateur
        st.session_state.chat_history.append({"role": "user", "content": user_input})
        
        # Appeler l'API pour modifier les critères
        try:
            system_prompt = """Tu es un assistant aidant à définir les critères d'évaluation pour des documents académiques.
            L'utilisateur peut demander à modifier, ajouter ou supprimer des critères.
            Réponds de manière concise et propose les nouveaux critères d'évaluation en format structuré."""
            
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Critères actuels:\n{st.session_state.criteria}\n\nDemande: {user_input}"}
                ],
                max_tokens=1000
            )
            
            assistant_response = response.choices[0].message.content
            st.session_state.chat_history.append({"role": "assistant", "content": assistant_response})
            
            # Mettre à jour les critères et sauvegarder
            st.session_state.criteria = assistant_response
            save_criteria(assistant_response)
            
        except Exception as e:
            st.error(f"Erreur API: {str(e)}")
        
        st.rerun()

with col2:
    st.subheader("📋 Critères Actuels")
    
    # Éditeur de texte pour les critères
    updated_criteria = st.text_area(
        "Critères d'évaluation:",
        value=st.session_state.criteria,
        height=400
    )
    
    # Mettre à jour les critères s'ils sont modifiés
    if updated_criteria != st.session_state.criteria:
        st.session_state.criteria = updated_criteria
        save_criteria(updated_criteria)  # Sauvegarder automatiquement
        st.success("✅ Critères sauvegardés!")

# Bouton d'évaluation
st.markdown("---")
st.subheader("🚀 Lancer l'Évaluation")

col_button, col_info = st.columns([1, 2])

with col_button:
    if st.button("▶️ Évaluer les Documents", use_container_width=True, type="primary"):
        st.session_state.evaluate_triggered = True

if "evaluate_triggered" in st.session_state and st.session_state.evaluate_triggered:
    # Lancer l'évaluation
    evaluate_documents(st.session_state.criteria)
    st.session_state.evaluate_triggered = False

# Résultats
if st.session_state.results:
    st.markdown("---")
    st.subheader("📊 Résultats de l'Évaluation")
    
    # Afficher le tableau des résultats
    df_results = pd.DataFrame(st.session_state.results)
    
    # Afficher les statistiques
    col_stat1, col_stat2, col_stat3, col_stat4 = st.columns(4)
    
    notes = [r['note'] for r in st.session_state.results if r['note'] is not None]
    
    with col_stat1:
        st.metric("Évalués", len([r for r in st.session_state.results if r['note'] is not None]))
    
    with col_stat2:
        st.metric("Moyenne", f"{sum(notes) / len(notes):.2f}/20" if notes else "N/A")
    
    with col_stat3:
        st.metric("Meilleure", f"{max(notes)}/20" if notes else "N/A")
    
    with col_stat4:
        st.metric("Pire", f"{min(notes)}/20" if notes else "N/A")
    
    # Tableau complet
    st.dataframe(df_results, use_container_width=True)
    
    # Télécharger Excel
    if st.button("💾 Télécharger en Excel"):
        filename = "evaluations_results.xlsx"
        df_results[["Prénom", "Nom", "Note /20"]].to_excel(filename, index=False)
        
        with open(filename, "rb") as f:
            st.download_button(
                label="📥 Fichier Excel",
                data=f.read(),
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
