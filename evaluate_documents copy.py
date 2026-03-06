import os
import re
import time
import base64
from openai import OpenAI
from dotenv import load_dotenv
from tqdm import tqdm
import pandas as pd
from pathlib import Path
from pptx import Presentation
try:
    from PyPDF2 import PdfReader
except:
    PdfReader = None

# Charger les variables d'environnement
load_dotenv("API.env")

# Initialiser le client OpenAI
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# DEFAULT EVALUATION PROMPT (peut être remplacé via l'app)
DEFAULT_EVALUATION_PROMPT = """Veuillez évaluer ce document académique sur la base des 5 critères suivants et attribuer une note totale sur 20 points (4 points par critère).

Les critères sont:
1. Plan adopté (4 pts) - Cohérence et organisation logique du contenu
2. Design des slides (4 pts) - Esthétique, lisibilité et qualité de la présentation visuelle
3. Référencement (4 pts) - Utilisation correcte des normes (APA ou autre) pour les sources
4. Existence de chiffres, tableaux ou graphiques (4 pts) - Présence et pertinence des éléments visuels de données
5. Richesse des informations (4 pts) - Profondeur du contenu, vocabulaire technique approprié et attesté par les références

Veuillez fournir:
- Une analyse succincte pour chaque critère (1-2 lignes max)
- La note obtenue pour chaque critère (0-4)
- La note totale sur 20
- Format de réponse strict:

CRITÈRE 1 - PLAN: [note]/4
[analyse succincte]

CRITÈRE 2 - DESIGN: [note]/4
[analyse succincte]

CRITÈRE 3 - RÉFÉRENCEMENT: [note]/4
[analyse succincte]

CRITÈRE 4 - DONNÉES VISUELLES: [note]/4
[analyse succincte]

CRITÈRE 5 - RICHESSE: [note]/4
[analyse succincte]

NOTE TOTALE: [note]/20
"""

def parse_filename(filename):
    """
    Extrait le prénom et nom du format: "DESCRIPTION - FirstName LASTNAME.extension"
    
    Args:
        filename (str): Nom du fichier
        
    Returns:
        tuple: (prénom, nom) ou (None, None) si extraction impossible
    """
    try:
        # Supprimer l'extension
        name_without_ext = Path(filename).stem
        
        # Chercher le tiret séparant la description du nom
        if " - " in name_without_ext:
            parts = name_without_ext.split(" - ")
            if len(parts) >= 2:
                # Prendre tout après le dernier tiret
                names = " - ".join(parts[1:]).strip()
                
                # Diviser par espaces pour obtenir prénom et nom
                name_parts = names.split()
                if len(name_parts) >= 2:
                    first_name = name_parts[0]
                    last_name = " ".join(name_parts[1:])
                    return first_name, last_name
                elif len(name_parts) == 1:
                    return name_parts[0], ""
    except Exception as e:
        print(f"Erreur lors du parsing du nom: {e}")
    
    return None, None

def extract_note(response_text):
    """
    Extrait la note totale du texte de réponse avec plusieurs patterns possibles.
    
    Args:
        response_text (str): Texte de réponse de l'API
        
    Returns:
        float: Note totale sur 20, ou None si non trouvée
    """
    try:
        # Pattern 1: "NOTE TOTALE: X/20"
        match = re.search(r'NOTE\s+TOTALE\s*:\s*(\d+(?:[.,]\d+)?)\s*/\s*20', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 2: "NOTE TOTALE: X" (sans /20)
        match = re.search(r'NOTE\s+TOTALE\s*:\s*(\d+(?:[.,]\d+)?)', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 3: Chercher juste "/20" suivi d'un nombre
        match = re.search(r'(\d+(?:[.,]\d+)?)\s*/\s*20', response_text)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
        # Pattern 4: Chercher "note:" ou "note =" suivi d'un nombre
        match = re.search(r'note\s*[:=]\s*(\d+(?:[.,]\d+)?)', response_text, re.IGNORECASE)
        if match:
            note = match.group(1).replace(',', '.')
            return float(note)
        
    except Exception as e:
        pass
    
    return None

def evaluate_document(file_id, file_name, evaluation_prompt=None):
    """
    Évalue un document en utilisant l'API OpenAI avec GPT-4 Vision.
    
    Args:
        file_id (str): ID du fichier uploadé
        file_name (str): Nom du fichier pour extraction du nom étudiant
        evaluation_prompt (str): Prompt d'évaluation (utilise le prompt par défaut si None)
        
    Returns:
        dict: Résultats de l'évaluation
    """
    # Utiliser le prompt par défaut si aucun n'est fourni
    if evaluation_prompt is None:
        evaluation_prompt = DEFAULT_EVALUATION_PROMPT
    
    first_name, last_name = parse_filename(file_name)
    
    try:
        # Chemins possibles du fichier uploadé
        documents_folder = r"C:\Users\hp\Downloads\Documents a analyser"
        file_path = os.path.join(documents_folder, file_name)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Fichier non trouvé: {file_path}")
        
        # Préparer le contenu pour l'API
        content = [
            {
                "type": "text",
                "text": evaluation_prompt
            }
        ]
        
        # Traiter le fichier selon son type
        file_ext = Path(file_name).suffix.lower()
        
        if file_ext == ".pdf":
            # Extraire le texte du PDF
            try:
                if PdfReader is None:
                    raise ImportError("PyPDF2 non installé")
                
                pdf_reader = PdfReader(file_path)
                text_content = "Contenu du PDF:\n"
                
                # Lire les premières 5 pages
                max_pages = min(5, len(pdf_reader.pages))
                for page_num in range(max_pages):
                    try:
                        page = pdf_reader.pages[page_num]
                        text_content += f"\n--- Page {page_num + 1} ---\n"
                        text_content += page.extract_text() + "\n"
                    except:
                        pass
                
                content.append({
                    "type": "text",
                    "text": text_content[:3000]  # Limiter la taille
                })
            except Exception as pdf_error:
                # En cas d'erreur, envoyer un message alternatif
                content.append({
                    "type": "text",
                    "text": f"Document PDF: {file_name}. Veuillez analyser ce document académique sur la base des critères fournis."
                })
        
        elif file_ext == ".pptx":
            # Extraire le texte de la présentation
            try:
                prs = Presentation(file_path)
                text_content = "Contenu de la présentation:\n"
                slide_count = min(5, len(prs.slides))  # Premières 5 slides
                
                for slide_idx in range(slide_count):
                    try:
                        slide = prs.slides[slide_idx]
                        text_content += f"\n--- Slide {slide_idx + 1} ---\n"
                        
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, "text") and shape.text:
                                    text_content += shape.text + "\n"
                            except:
                                pass
                    except:
                        continue
                
                # Ajouter le texte au contenu
                content.append({
                    "type": "text",
                    "text": text_content[:3000]  # Limiter la taille
                })
            except Exception as pptx_error:
                # En cas d'erreur PPTX, envoyer un message alternatif
                content.append({
                    "type": "text",
                    "text": f"Document PPTX: {file_name}. Veuillez analyser ce document de présentation académique."
                })
        
        # Appeler l'API GPT-4 Vision
        response = client.chat.completions.create(
            model="gpt-4o",
            max_tokens=1024,
            messages=[
                {
                    "role": "user",
                    "content": content
                }
            ]
        )
        
        evaluation = response.choices[0].message.content
        note = extract_note(evaluation)
        
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": note,
            "evaluation": evaluation,
            "status": "✅ Succès" if note is not None else "⚠️ Note non trouvée"
        }
    
    except FileNotFoundError as e:
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": None,
            "evaluation": str(e),
            "status": f"❌ Fichier non trouvé"
        }
    
    except Exception as e:
        return {
            "file_name": file_name,
            "first_name": first_name,
            "last_name": last_name,
            "note": None,
            "evaluation": str(e),
            "status": f"❌ Erreur: {str(e)[:50]}"
        }

def read_upload_results(filename="upload_results.txt"):
    """
    Lit le fichier de résultats d'upload pour extraire les IDs de fichiers.
    
    Args:
        filename (str): Nom du fichier de résultats
        
    Returns:
        list: Liste de tuples (file_id, file_name)
    """
    results = []
    
    try:
        with open(filename, "r", encoding="utf-8") as f:
            lines = f.readlines()
        
        current_file = None
        for line in lines:
            if line.startswith("Fichier:"):
                current_file = line.replace("Fichier:", "").strip()
            elif line.startswith("ID:") and current_file:
                file_id = line.replace("ID:", "").strip()
                results.append((file_id, current_file))
                current_file = None
        
        return results
    
    except FileNotFoundError:
        print(f"❌ Fichier {filename} non trouvé")
        return []
    except Exception as e:
        print(f"❌ Erreur lors de la lecture du fichier: {e}")
        return []

def save_to_excel(results, output_file="evaluations.xlsx"):
    """
    Sauvegarde les résultats d'évaluation dans un fichier Excel.
    
    Args:
        results (list): Liste des résultats d'évaluation
        output_file (str): Nom du fichier Excel à créer
    """
    try:
        # Préparer les données
        data = {
            "Prénom": [r.get("first_name", "") for r in results],
            "Nom": [r.get("last_name", "") for r in results],
            "Note /20": [r.get("note", "") for r in results],
            "Statut": [r.get("status", "") for r in results]
        }
        
        # Créer un DataFrame
        df = pd.DataFrame(data)
        
        # Sauvegarder en Excel
        df.to_excel(output_file, index=False, engine='openpyxl')
        
        print(f"\n💾 Résultats sauvegardés dans: {output_file}")
        return True
    
    except ImportError:
        print("⚠️  pandas ou openpyxl n'est pas installé")
        return False
    except Exception as e:
        print(f"❌ Erreur lors de la sauvegarde: {e}")
        return False

def evaluate_batch(evaluation_prompt=None):
    """
    Évalue tous les documents en batch.
    
    Args:
        evaluation_prompt (str): Prompt d'évaluation personnalisé (optionnel)
    """
    # Utiliser le prompt par défaut si aucun n'est fourni
    if evaluation_prompt is None:
        evaluation_prompt = DEFAULT_EVALUATION_PROMPT
    
    print("=" * 60)
    print("🔍 ÉVALUATION EN BATCH DES DOCUMENTS")
    print("=" * 60)
    
    # Lire les fichiers uploadés
    files_to_evaluate = read_upload_results()
    
    if not files_to_evaluate:
        print("⚠️  Aucun fichier à évaluer. Vérifiez upload_results.txt")
        return
    
    print(f"\n📊 Fichiers à évaluer: {len(files_to_evaluate)}\n")
    
    # Évaluer chaque document
    results = []
    for file_id, file_name in tqdm(files_to_evaluate, desc="Évaluation en cours", unit="document"):
        result = evaluate_document(file_id, file_name, evaluation_prompt)
        results.append(result)
        
        # Afficher le résultat
        print(f"{result['status']} - {result['file_name']}")
        if result['note'] is not None:
            print(f"   Note: {result['note']}/20")
    
    # Résumé
    print("\n" + "=" * 60)
    print("📊 RÉSUMÉ")
    print("=" * 60)
    successful = sum(1 for r in results if r['note'] is not None)
    print(f"✅ Évalués avec succès: {successful}/{len(results)}")
    
    # Statistiques
    notes = [r['note'] for r in results if r['note'] is not None]
    if notes:
        print(f"📈 Moyenne: {sum(notes) / len(notes):.2f}/20")
        print(f"🔝 Meilleure note: {max(notes)}/20")
        print(f"🔻 Pire note: {min(notes)}/20")
    
    # Sauvegarder en Excel
    save_to_excel(results)

if __name__ == "__main__":
    evaluate_batch()
